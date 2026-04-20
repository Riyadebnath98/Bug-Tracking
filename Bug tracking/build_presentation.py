"""
Generates academic PPT for Bug Tracking System mini project.
Run: python build_presentation.py
Output: RiyaDebnath_SoftwareEngineering_BugTrackingSystem.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


def set_title(slide, text: str) -> None:
    slide.shapes.title.text = text
    tf = slide.shapes.title.text_frame
    for p in tf.paragraphs:
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)


def add_bullets(prs, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    set_title(slide, title)
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            body.text = b
        else:
            p = body.add_paragraph()
            p.text = b
            p.level = 0
        body.paragraphs[-1].font.size = Pt(17)
        body.paragraphs[-1].space_after = Pt(4)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- 1. Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "BUG TRACKING SYSTEM"
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    sub = slide.placeholders[1]
    tf = sub.text_frame
    tf.text = "Software Engineering Mini Project"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Presented by: Riya Debnath"
    p.font.size = Pt(20)
    p = tf.add_paragraph()
    p.text = "Guide: Mr. Likith"
    p.font.size = Pt(20)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Department of Computer Science and Engineering"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "(Update institution name on this slide if required)"
    p.font.size = Pt(14)
    p.font.italic = True

    # --- 2. Contents ---
    contents = [
        "1. Title Slide   2. Contents / Index   3. Abstract",
        "4. Introduction / Background   5. Objective of the Study   6. Problem Definition",
        "7. Scope of the Study   8. Literature Review   9. Methodology / Approach",
        "10. Experimental Setup or Data Collection   11. Analysis and Interpretation",
        "12. Results or Findings   13. Consolidation Table / Summary Table   14. Discussion",
        "15. Proposed Solutions or Proposed System   16. Feasibility Study",
        "17. Design / Workflow / Process Design   18. Diagrams",
        "19. Methods and Techniques Used   20. Progress in Project",
        "21. Outcome / Results   22. Conclusion   23. References",
    ]
    add_bullets(prs, "2. Contents / Index", contents)

    # --- 3. Abstract ---
    add_bullets(
        prs,
        "3. Abstract",
        [
            "Software teams need a reliable way to record defects so issues are tracked from report to closure.",
            "This project implements a web-based Bug Tracking System for logging, assigning, and monitoring bugs.",
            "The system supports user roles—Developer, Tester, and Project Manager—and stores data in SQLite.",
            "Users can create, read, update, and delete bug records; search and filter by text, status, priority, and assignee.",
            "An activity log provides recent notifications when bugs are created, edited, or change status.",
            "The application runs locally for demonstration and can be extended for broader deployment.",
        ],
    )

    # --- 4. Introduction ---
    add_bullets(
        prs,
        "4. Introduction / Background",
        [
            "Bug tracking is the practice of capturing software defects, prioritizing them, and following them until resolved.",
            "Without a dedicated tool, teams often rely on informal channels, which leads to lost context and duplicate work.",
            "Structured tracking improves communication between testers who report issues and developers who fix them.",
            "Project managers benefit from visibility into status and workload across the team.",
            "This project applies standard software engineering ideas to a small, practical web application.",
        ],
    )

    # --- 5. Objective ---
    add_bullets(
        prs,
        "5. Objective of the Study",
        [
            "To design and implement a bug tracking web application with secure login and session-based access.",
            "To maintain bug records with description, priority, status, creator, and assignee.",
            "To provide full CRUD operations and quick retrieval through search and filters.",
            "To record notable actions in an activity log for traceability and simple in-app notifications.",
            "To demonstrate local deployment using Python, Flask, and SQLite.",
        ],
    )

    # --- 6. Problem Definition ---
    add_bullets(
        prs,
        "6. Problem Definition",
        [
            "Problem: Ad-hoc bug reporting is inconsistent and does not scale when multiple people work on one product.",
            "Need: A single repository where each bug has clear attributes and a visible lifecycle (e.g., Open → In Progress → Resolved).",
            "Need: Ability to assign ownership and filter work so developers and managers can focus on relevant items.",
            "Scope of this work: A prototype suitable for coursework and demo, not a full enterprise product.",
            "Known limitations: Plain-text passwords in the demo database; role labels exist but fine-grained permission rules can be added later.",
        ],
    )

    # --- 7. Scope ---
    add_bullets(
        prs,
        "7. Scope of the Study",
        [
            "In scope: User authentication, dashboard listing, bug create/edit/delete, status updates, search and multi-criteria filters.",
            "In scope: SQLite persistence, schema for users, bugs, and activity_log; seeded demo users and sample bugs.",
            "In scope: Server-rendered HTML pages and CSS styling for a clean, usable interface.",
            "Out of scope: Email or SMS alerts, mobile native apps, multi-tenant cloud hosting, and advanced analytics dashboards.",
            "Out of scope: Production-grade password hashing and OAuth—recommended as future enhancements.",
        ],
    )

    # --- 8. Literature Review ---
    add_bullets(
        prs,
        "8. Literature Review (Optional but Recommended)",
        [
            "Industry issue trackers (e.g., Jira, Azure DevOps, Bugzilla) show common patterns: work items, workflows, and audit history.",
            "Software engineering texts describe defect lifecycle models and the value of traceability from requirement to fix.",
            "Lightweight web stacks (microframeworks plus embedded databases) are widely used for teaching and rapid prototypes.",
            "Takeaway: This project adopts familiar bug attributes and status flow while keeping the stack minimal for learning.",
            "Note: Expand this section with 2–3 papers or official docs you have read; paraphrase in your own words for submission.",
        ],
    )

    # --- 9. Methodology ---
    add_bullets(
        prs,
        "9. Methodology / Approach",
        [
            "Requirements were gathered from the mini-project specification and guide expectations.",
            "Data modeling: entities for users, bugs, and activity messages with keys and allowed values for priority and status.",
            "Application structure: Flask application factory, separate blueprints for authentication and bug routes.",
            "Implementation: SQL queries with bound parameters; Jinja templates for views; CSS for layout and readability.",
            "Validation: Manual testing of login, CRUD, filters, and activity log entries after each milestone.",
        ],
    )

    # --- 10. Experimental Setup ---
    add_bullets(
        prs,
        "10. Experimental Setup or Data Collection (Domain-specific)",
        [
            "Hardware: A standard PC running Windows; tests performed in a desktop web browser.",
            "Software: Python 3.x, Flask 3.1.0, SQLite via the standard library; project folder as the working directory.",
            "Deployment: python app.py; application served at http://localhost:5000; login page at /login.",
            "Test accounts: manager/manager123 (Project Manager), tester1/tester123 (Tester), dev1/dev123 (Developer).",
            "Data: bugtracker.db created automatically; sample bugs seeded on first run for demonstration.",
        ],
    )

    # --- 11. Analysis ---
    add_bullets(
        prs,
        "11. Analysis and Interpretation",
        [
            "Each stated requirement maps to a module: auth.py for login/logout, bugs.py for CRUD and filters, db.py for schema and seeding.",
            "Search combines title and description; filters restrict by allowed status and priority values to avoid invalid input.",
            "Activity log entries are written on create, detail update, and status change—supporting audit and the notification panel.",
            "The design favors clarity over scale: suitable for demos; high concurrency would need a different deployment model.",
            "User experience is kept simple so testers and developers can complete tasks with minimal training.",
        ],
    )

    # --- 12. Results ---
    add_bullets(
        prs,
        "12. Results or Findings",
        [
            "Login and session protection: unauthenticated users are redirected to the login page.",
            "Dashboard shows the bug list with creator and assignee names, plus recent activity messages.",
            "Users can submit new bugs with title, description, priority, and optional assignee.",
            "Edit page updates details; inline status control moves bugs through Open, In Progress, and Resolved.",
            "Delete removes a bug and its related activity rows; search and filters narrow the list as expected.",
            "Add screenshots of login, dashboard, and edit screens to this deck and describe each in the speaker notes or caption text.",
        ],
    )

    # --- 13. Consolidation Table ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "13. Consolidation Table / Summary Table"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    rows, cols = 9, 3
    left, top, width, height = Inches(0.5), Inches(1.2), Inches(9), Inches(5.5)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(3.8)
    table.columns[2].width = Inches(2.0)

    headers = ("Requirement / Feature", "Implementation summary", "Status")
    data = [
        ("Log, track, manage bugs", "Bug records with lifecycle and assignee", "Done"),
        ("Roles: Developer, Tester, PM", "User table with role field; demo accounts", "Done"),
        ("Bug fields: id, description, priority, status, assignee, created_by", "SQLite schema + title and timestamps", "Done"),
        ("Full CRUD", "Create, list, edit, delete routes", "Done"),
        ("Search and filter", "Query string filters for text, status, priority, assignee", "Done"),
        ("Notifications", "activity_log table; last 8 events on dashboard", "Done"),
        ("Sample data", "Seeded users and bugs on first DB init", "Done"),
        ("Local deployment", "Flask dev server on port 5000", "Done"),
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(14)
    for r, row in enumerate(data, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)

    # --- 14. Discussion ---
    add_bullets(
        prs,
        "14. Discussion",
        [
            "The prototype meets the core academic goals: structured data, roles, CRUD, and traceability via the activity log.",
            "SQLite keeps setup simple; migrating to PostgreSQL or MySQL would be natural if the user base grows.",
            "Session-based login is appropriate for coursework; production systems would add password hashing and HTTPS.",
            "Role labels could drive future rules—for example, only testers create bugs or only assignees change status.",
            "Overall, the project balances learning objectives with a coherent end-to-end user flow.",
        ],
    )

    # --- 15. Proposed System ---
    add_bullets(
        prs,
        "15. Proposed Solutions or Proposed System",
        [
            "Proposed system: A browser-based Bug Tracking System built with Flask and SQLite.",
            "Users authenticate once per session; the dashboard centralizes all bug operations.",
            "Bugs carry priority (Low to Critical) and status (Open, In Progress, Resolved) to mirror common team practice.",
            "The activity log acts as a lightweight notification stream without external messaging services.",
            "The codebase is organized into db, auth, bugs, templates, and static assets for maintainability.",
        ],
    )

    # --- 16. Feasibility ---
    add_bullets(
        prs,
        "16. Feasibility Study (Technical / Business Domains Only)",
        [
            "Technical feasibility: Flask and SQLite are stable, well-documented, and sufficient for a local prototype.",
            "Economic feasibility: Open-source tools; no licensing cost for the core stack used here.",
            "Operational feasibility: Single-command startup; easy to demonstrate on a laptop for internal evaluation.",
            "Schedule feasibility: Features align with a typical mini-project timeline when developed incrementally.",
            "Risks: Security and scalability are limited by design; both are acceptable for academic scope with a clear upgrade path.",
        ],
    )

    # --- 17. Design / Workflow ---
    add_bullets(
        prs,
        "17. Design / Workflow / Process Design",
        [
            "Start → User opens login page → submits credentials → session created → redirect to dashboard.",
            "Dashboard → filter or search → open edit page to change details → save → activity log updated.",
            "Dashboard → change status via POST → timestamp updated → log message recorded.",
            "Create bug form → validate required fields → insert row → log creation event.",
            "Logout clears session and returns to login; database connections close per Flask app context.",
        ],
    )

    # --- 18. Diagrams (multiple slides, same section) ---
    add_bullets(
        prs,
        "18. Diagrams — High-Level Architecture",
        [
            "Browser (HTML/CSS) ↔ Flask server (routes, session, Jinja rendering).",
            "Flask application ↔ SQLite file (bugtracker.db): users, bugs, activity_log tables.",
            "Static assets (styles.css) served for consistent layout and readability.",
            "Replace or augment this slide with a labeled architecture diagram drawn in your preferred tool.",
        ],
    )
    add_bullets(
        prs,
        "18. Diagrams — Use Case Overview",
        [
            "Actors: Authenticated user (roles: Developer, Tester, Project Manager).",
            "Use cases: Login, Logout, View dashboard, Search/filter bugs, Create bug, Edit bug, Update status, Delete bug, View notifications.",
            "All bug operations require an active session (login_required decorator).",
            "Draw a standard use case diagram with actors and ovals; label associations clearly.",
        ],
    )
    add_bullets(
        prs,
        "18. Diagrams — Data Flow (Conceptual)",
        [
            "Level 0: User provides inputs (forms) → System processes requests → System reads/writes database → System returns HTML pages.",
            "Level 1 processes: Authentication, Bug management (CRUD), Reporting (list + filters), Activity logging.",
            "Prepare a DFD with numbered processes, data stores (D1 Users, D2 Bugs, D3 Activity), and labeled data flows.",
        ],
    )
    add_bullets(
        prs,
        "18. Diagrams — Entity Relationship (Conceptual)",
        [
            "User (id, username, password, role) — creates — Bug; User may be assignee of Bug (optional).",
            "Bug (id, title, description, priority, status, assignee_id, created_by, timestamps).",
            "Activity_log (id, bug_id, message, created_at) — many rows per Bug.",
            "Draw ERD with cardinality: User 1—* Bug (created_by); User 1—* Bug (assignee); Bug 1—* Activity_log.",
        ],
    )

    # --- 19. Methods ---
    add_bullets(
        prs,
        "19. Methods and Techniques Used",
        [
            "Programming language: Python.",
            "Web framework: Flask 3.1.0 with blueprints and application factory pattern.",
            "Database: SQLite; API via sqlite3 with Row factory for dict-like access.",
            "Presentation: Jinja2 HTML templates; CSS for styling.",
            "Authentication: Server-side session after form login; parameterized SQL to reduce injection risk.",
        ],
    )

    # --- 20. Progress ---
    add_bullets(
        prs,
        "20. Progress in Project",
        [
            "Phase 1: Requirement understanding and environment setup (Python, Flask, virtual environment).",
            "Phase 2: Database schema, migrations via init script, seed users and sample bugs.",
            "Phase 3: Login, logout, and protected routes.",
            "Phase 4: Bug CRUD, status updates, search and filter, activity log integration.",
            "Phase 5: UI polish, README documentation, presentation preparation—adjust dates to match your actual plan.",
        ],
    )

    # --- 21. Outcome ---
    add_bullets(
        prs,
        "21. Outcome / Results",
        [
            "Deliverable: A working Bug Tracking System prototype aligned with the mini-project checklist.",
            "Demonstrable flows for multiple roles using seeded accounts.",
            "Clear separation of configuration, data access, authentication, and bug features in source files.",
            "Foundation for future work: password hashing, role-based permissions, deployment to a production server.",
        ],
    )

    # --- 22. Conclusion ---
    add_bullets(
        prs,
        "22. Conclusion",
        [
            "The project addresses the need for structured defect tracking using a lightweight web stack.",
            "Core objectives—CRUD, roles, persistence, search, filters, and activity history—are implemented successfully.",
            "The solution is appropriate for academic evaluation and small-team demos.",
            "Lessons learned include practical use of Flask, SQL design, and end-to-end feature integration.",
            "Thank you. Questions are welcome.",
        ],
    )

    # --- 23. References ---
    add_bullets(
        prs,
        "23. References",
        [
            "Flask Documentation. https://flask.palletsprojects.com/ (accessed as per your study period).",
            "SQLite Documentation. https://www.sqlite.org/docs.html",
            "Pallets Project — Jinja2. https://jinja.palletsprojects.com/",
            "Sommerville, I. Software Engineering (or your prescribed textbook) — chapters on testing and configuration management.",
            "Add any journal or conference papers you cited in Section 8, using your department’s citation style (IEEE / APA).",
        ],
    )

    out_dir = Path(__file__).resolve().parent
    out_path = out_dir / "RiyaDebnath_SoftwareEngineering_BugTrackingSystem.pptx"
    prs.save(str(out_path))
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
