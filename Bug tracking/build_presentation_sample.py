"""
Builds PPT strictly in the sample-PDF section order.
Run: python build_presentation_sample.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt


def set_title(slide, title_text):
    slide.shapes.title.text = title_text
    t = slide.shapes.title.text_frame.paragraphs[0]
    t.font.size = Pt(30)
    t.font.bold = True
    t.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, title)
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            tf.text = b
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(5)


def add_consolidation_table(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.8))
    tf = box.text_frame
    tf.text = "7. CONSOLIDATION TABLE"
    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    rows, cols = 8, 3
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(9.0), Inches(5.8)).table
    table.columns[0].width = Inches(3.1)
    table.columns[1].width = Inches(4.1)
    table.columns[2].width = Inches(1.8)

    headers = ["Requirement", "Implementation in Bug Tracking System", "Status"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        hp = cell.text_frame.paragraphs[0]
        hp.font.bold = True
        hp.font.size = Pt(13)

    rows_data = [
        ("Roles", "Developer, Tester, Project Manager in users table", "Done"),
        ("Bug lifecycle", "Open, In Progress, Resolved workflow", "Done"),
        ("CRUD", "Create, view, edit, delete bug records", "Done"),
        ("Search/filter", "Text + status + priority + assignee filters", "Done"),
        ("Notifications", "activity_log entries shown on dashboard", "Done"),
        ("Data storage", "SQLite schema in db.py", "Done"),
        ("Local deployment", "Runs on localhost:5000", "Done"),
    ]
    for r, row in enumerate(rows_data, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cp = cell.text_frame.paragraphs[0]
            cp.font.size = Pt(12)


def add_box(slide, x, y, w, h, text, font_size=13, fill=(230, 240, 255)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    tf = shape.text_frame
    tf.clear()
    tf.text = text
    for p in tf.paragraphs:
        p.font.size = Pt(font_size)
        p.font.bold = False
    return shape


def connect(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    line.line.width = Pt(1.5)
    return line


def add_architecture_diagram_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.8))
    tf = box.text_frame
    tf.text = "12. ARCHITECTURE DIAGRAM"
    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    # Main 3-tier blocks
    add_box(slide, 0.7, 2.2, 2.2, 1.3, "Client Browser\n(Login/Dashboard)", 12)
    add_box(slide, 3.9, 2.0, 2.3, 1.7, "Flask App\n(auth.py / bugs.py)\nBusiness Logic", 12)
    add_box(slide, 7.1, 2.2, 2.0, 1.3, "SQLite DB\nusers, bugs,\nactivity_log", 12)

    connect(slide, 2.9, 2.85, 3.9, 2.85)
    connect(slide, 3.9, 3.05, 2.9, 3.05)
    connect(slide, 6.2, 2.85, 7.1, 2.85)
    connect(slide, 7.1, 3.05, 6.2, 3.05)

    # Supporting components
    add_box(slide, 3.9, 4.4, 2.3, 1.0, "Jinja Templates\nHTML Rendering", 11, fill=(245, 245, 245))
    add_box(slide, 0.7, 4.4, 2.2, 1.0, "Static CSS\n(styles.css)", 11, fill=(245, 245, 245))
    connect(slide, 5.0, 3.7, 5.0, 4.4)
    connect(slide, 3.9, 4.9, 2.9, 4.9)


def add_use_case_diagram_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.8))
    tf = box.text_frame
    tf.text = "13. USE CASE DIAGRAM"
    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    # Actors
    add_box(slide, 0.4, 2.0, 1.6, 0.8, "Tester", 12, fill=(255, 245, 230))
    add_box(slide, 0.4, 3.2, 1.6, 0.8, "Developer", 12, fill=(255, 245, 230))
    add_box(slide, 0.4, 4.4, 1.6, 0.8, "Project Manager", 12, fill=(255, 245, 230))

    # System boundary
    boundary = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.3), Inches(1.6), Inches(7.2), Inches(4.3))
    boundary.fill.background()
    boundary.line.color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    btf = boundary.text_frame
    btf.text = "Bug Tracking System"
    btf.paragraphs[0].font.bold = True
    btf.paragraphs[0].font.size = Pt(12)

    # Use cases as ovals
    uc1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.8), Inches(2.1), Inches(1.8), Inches(0.7))
    uc1.text_frame.text = "Login"
    uc2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.0), Inches(2.1), Inches(2.2), Inches(0.7))
    uc2.text_frame.text = "Create / Edit Bug"
    uc3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(2.1), Inches(1.6), Inches(0.7))
    uc3.text_frame.text = "Delete Bug"
    uc4 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.2), Inches(3.4), Inches(2.4), Inches(0.7))
    uc4.text_frame.text = "Update Status"
    uc5 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.1), Inches(3.4), Inches(2.8), Inches(0.7))
    uc5.text_frame.text = "Search / Filter"
    uc6 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.8), Inches(4.7), Inches(2.0), Inches(0.7))
    uc6.text_frame.text = "View Activity"

    for uc in [uc1, uc2, uc3, uc4, uc5, uc6]:
        uc.fill.solid()
        uc.fill.fore_color.rgb = RGBColor(240, 248, 255)
        uc.line.color.rgb = RGBColor(0x1A, 0x23, 0x7E)
        uc.text_frame.paragraphs[0].font.size = Pt(11)

    # Associations
    connect(slide, 2.0, 2.4, 2.8, 2.4)
    connect(slide, 2.0, 2.4, 5.0, 2.4)
    connect(slide, 2.0, 3.6, 5.0, 2.4)
    connect(slide, 2.0, 3.6, 3.2, 3.7)
    connect(slide, 2.0, 4.8, 7.5, 2.4)
    connect(slide, 2.0, 4.8, 6.1, 3.7)
    connect(slide, 2.0, 4.8, 4.8, 5.0)


def add_dfd_diagram_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.8))
    tf = box.text_frame
    tf.text = "14. DATA FLOW DIAGRAM"
    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    add_box(slide, 0.4, 2.8, 1.8, 0.9, "User", 12, fill=(255, 245, 230))
    add_box(slide, 3.0, 2.6, 2.4, 1.2, "Process:\nBug Tracking\nSystem", 12)
    add_box(slide, 6.6, 1.9, 2.4, 1.0, "D1: Users", 12, fill=(245, 245, 245))
    add_box(slide, 6.6, 3.1, 2.4, 1.0, "D2: Bugs", 12, fill=(245, 245, 245))
    add_box(slide, 6.6, 4.3, 2.4, 1.0, "D3: Activity Log", 12, fill=(245, 245, 245))

    connect(slide, 2.2, 3.1, 3.0, 3.1)   # input
    connect(slide, 5.4, 2.95, 6.6, 2.4)  # read users
    connect(slide, 5.4, 3.1, 6.6, 3.6)   # read/write bugs
    connect(slide, 5.4, 3.25, 6.6, 4.8)  # write log
    connect(slide, 3.0, 3.35, 2.2, 3.35)  # output

    # Data labels
    lbl1 = slide.shapes.add_textbox(Inches(2.25), Inches(2.72), Inches(0.8), Inches(0.3))
    lbl1.text_frame.text = "Input"
    lbl1.text_frame.paragraphs[0].font.size = Pt(10)
    lbl2 = slide.shapes.add_textbox(Inches(2.25), Inches(3.38), Inches(1.2), Inches(0.3))
    lbl2.text_frame.text = "Output"
    lbl2.text_frame.paragraphs[0].font.size = Pt(10)


def add_erd_diagram_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.8))
    tf = box.text_frame
    tf.text = "15. ER DIAGRAM"
    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    user = add_box(
        slide,
        0.6,
        2.0,
        2.7,
        2.0,
        "USERS\n- id (PK)\n- username\n- password\n- role",
        11,
        fill=(255, 250, 235),
    )
    bug = add_box(
        slide,
        3.8,
        1.7,
        3.0,
        2.6,
        "BUGS\n- id (PK)\n- title\n- description\n- priority\n- status\n- assignee_id (FK)\n- created_by (FK)",
        10,
        fill=(235, 250, 255),
    )
    add_box(
        slide,
        7.2,
        2.0,
        2.2,
        2.0,
        "ACTIVITY_LOG\n- id (PK)\n- bug_id (FK)\n- message\n- created_at",
        10,
        fill=(245, 245, 245),
    )

    connect(slide, 3.3, 2.7, 3.8, 2.7)  # users->bugs
    connect(slide, 6.8, 2.9, 7.2, 2.9)  # bugs->activity

    c1 = slide.shapes.add_textbox(Inches(3.2), Inches(2.35), Inches(0.8), Inches(0.3))
    c1.text_frame.text = "1..*"
    c1.text_frame.paragraphs[0].font.size = Pt(10)
    c2 = slide.shapes.add_textbox(Inches(6.75), Inches(2.55), Inches(0.8), Inches(0.3))
    c2.text_frame.text = "1..*"
    c2.text_frame.paragraphs[0].font.size = Pt(10)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide like sample
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "BUG TRACKING SYSTEM"
    t = slide.shapes.title.text_frame.paragraphs[0]
    t.font.size = Pt(38)
    t.font.bold = True
    t.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)
    sub = slide.placeholders[1].text_frame
    sub.text = "Department of Computer Science and Engineering"
    p = sub.add_paragraph()
    p.text = ""
    p = sub.add_paragraph()
    p.text = "Presented by: RIYA DEBNATH"
    p.font.size = Pt(20)
    p = sub.add_paragraph()
    p.text = "Guide: Mr. Likith"
    p.font.size = Pt(20)
    p = sub.add_paragraph()
    p.text = "Final Review"
    p.font.size = Pt(16)

    # Sample-PDF exact contents headings/order
    add_bullet_slide(
        prs,
        "CONTENTS",
        [
            "1 ABSTRACT",
            "2 INTRODUCTION",
            "3 OBJECTIVE",
            "4 SCOPE",
            "5 PROBLEM DEFINITION",
            "6 LITERATURE SURVEY",
            "7 CONSOLIDATION TABLE",
            "8 REQUIREMENT SPECIFICATION",
            "9 PROPOSED SYSTEM",
            "10 FEASIBILITY STUDY",
            "11 DESIGN",
            "12 ARCHITECTURE DIAGRAM",
            "13 USE CASE DIAGRAM",
            "14 DATA FLOW DIAGRAM",
            "15 ER DIAGRAM",
            "16 METHODS AND TECHNIQUES",
            "17 PROGRESS IN PROJECT",
            "18 OUTCOME",
            "19 CONCLUSION",
            "20 REFERENCES",
        ],
    )

    add_bullet_slide(
        prs,
        "1. ABSTRACT",
        [
            "Software teams need a structured way to register and track defects.",
            "This project presents a web-based Bug Tracking System using Flask and SQLite.",
            "The system supports bug logging, assignment, priority management, and status updates.",
            "Users can search/filter issues and view recent activity notifications.",
            "The prototype is designed for local deployment and academic demonstration.",
        ],
    )

    add_bullet_slide(
        prs,
        "2. INTRODUCTION",
        [
            "Bug tracking is essential to maintain quality in software development.",
            "Manual methods (chat/spreadsheets) often cause duplicate work and missed defects.",
            "A centralized tracker improves communication between Tester, Developer, and Project Manager.",
            "This project implements that workflow in a lightweight web application.",
        ],
    )

    add_bullet_slide(
        prs,
        "3. OBJECTIVE",
        [
            "Build a working bug tracking system with role-based users.",
            "Support full bug lifecycle from creation to resolution.",
            "Enable quick bug analysis through search and filters.",
            "Maintain action history through activity logs.",
            "Provide a simple, demonstrable solution for mini-project evaluation.",
        ],
    )

    add_bullet_slide(
        prs,
        "4. SCOPE",
        [
            "Detect and record software issues reported by testers/users.",
            "Assign bugs to developers and monitor progress.",
            "Prioritize work using Low/Medium/High/Critical levels.",
            "Provide real-time dashboard visibility for project monitoring.",
            "Offer a maintainable base for future feature expansion.",
        ],
    )

    add_bullet_slide(
        prs,
        "5. PROBLEM DEFINITION",
        [
            "Software teams face delays when bugs are not tracked systematically.",
            "Important issue details are often missing in informal communication channels.",
            "Without proper ownership and status flow, bug resolution becomes inconsistent.",
            "This project solves the problem by storing standardized bug records in one system.",
        ],
    )

    add_bullet_slide(
        prs,
        "6. LITERATURE SURVEY",
        [
            "Common issue trackers follow shared principles: centralized bug records and workflow states.",
            "Research on defect management highlights traceability, prioritization, and accountability.",
            "Modern web applications often use modular architecture for maintainability.",
            "This project adopts these practices in a simplified academic implementation.",
        ],
    )

    add_consolidation_table(prs)

    add_bullet_slide(
        prs,
        "8. REQUIREMENT SPECIFICATION",
        [
            "Functional: login/logout, create/read/update/delete bugs, status changes.",
            "Functional: search by text and filter by status, priority, assignee.",
            "Functional: activity log for major bug events.",
            "Non-functional: simple UI, local execution, maintainable code structure.",
            "Tech stack: Python, Flask, SQLite, HTML templates, CSS.",
        ],
    )

    add_bullet_slide(
        prs,
        "9. PROPOSED SYSTEM",
        [
            "A Flask web application with blueprint-based modules.",
            "Authentication module controls session access.",
            "Bug module handles dashboard listing, CRUD operations, and status changes.",
            "SQLite database stores users, bugs, and activity logs.",
            "Templates render responsive pages for daily bug management tasks.",
        ],
    )

    add_bullet_slide(
        prs,
        "10. FEASIBILITY STUDY",
        [
            "Technical: The selected stack is stable and suitable for small systems.",
            "Economic: Open-source tools minimize implementation cost.",
            "Operational: Easy to run and demonstrate on a standard laptop.",
            "Schedule: The project can be completed within mini-project timelines.",
            "Risk: Advanced security/scalability features can be future enhancements.",
        ],
    )

    add_bullet_slide(
        prs,
        "11. DESIGN",
        [
            "Modular design with separate files for database, authentication, and bug routes.",
            "Entity structure: users, bugs, activity_log.",
            "Session-based access control for protected pages.",
            "UI design emphasizes clarity and quick bug operations.",
        ],
    )

    add_architecture_diagram_slide(prs)
    add_use_case_diagram_slide(prs)
    add_dfd_diagram_slide(prs)
    add_erd_diagram_slide(prs)

    add_bullet_slide(
        prs,
        "16. METHODS AND TECHNIQUES",
        [
            "Flask application factory and blueprints for modular development.",
            "SQLite relational modeling with constraints on role/priority/status.",
            "Parameterized SQL queries for safer database interactions.",
            "Session-based authentication and route protection.",
            "Server-side rendering using Jinja templates.",
        ],
    )

    add_bullet_slide(
        prs,
        "17. PROGRESS IN PROJECT",
        [
            "Phase 1: Requirement understanding and initial project setup.",
            "Phase 2: Database schema creation and sample data seeding.",
            "Phase 3: Authentication and dashboard implementation.",
            "Phase 4: CRUD, status updates, filters, and activity log features.",
            "Phase 5: Testing, documentation, and presentation preparation.",
        ],
    )

    add_bullet_slide(
        prs,
        "18. OUTCOME",
        [
            "A complete mini-project prototype for bug lifecycle management.",
            "All core features are available and testable on localhost.",
            "The system demonstrates practical software engineering workflow integration.",
            "The design is extendable for future institutional or industrial use.",
        ],
    )

    add_bullet_slide(
        prs,
        "19. CONCLUSION",
        [
            "The Bug Tracking System addresses real project coordination issues effectively.",
            "Core targets—tracking, assignment, prioritization, and monitoring—are achieved.",
            "The project provides a strong foundation for future secure and scalable enhancements.",
            "Thank you.",
        ],
    )

    add_bullet_slide(
        prs,
        "20. REFERENCES",
        [
            "Flask Official Documentation: https://flask.palletsprojects.com/",
            "SQLite Official Documentation: https://www.sqlite.org/docs.html",
            "Jinja Documentation: https://jinja.palletsprojects.com/",
            "Software Engineering notes/textbook used for mini-project preparation.",
            "Add any additional faculty-approved references used in your report.",
        ],
    )

    output = Path(__file__).resolve().parent / "RiyaDebnath_BugTrackingSystem_SampleFormat_v2.pptx"
    prs.save(str(output))
    print(f"Saved: {output}")


if __name__ == "__main__":
    main()
